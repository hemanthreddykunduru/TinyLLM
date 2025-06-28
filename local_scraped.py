import os
from pathlib import Path
from docx import Document
from pptx import Presentation
import PyPDF2
import pandas as pd
import textract

def read_text_file(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def read_docx_file(filepath):
    doc = Document(filepath)
    return '\n'.join([para.text for para in doc.paragraphs])

def read_pdf_file(filepath):
    text = ""
    with open(filepath, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def read_pptx_file(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_spreadsheet(filepath):
    try:
        df = pd.read_excel(filepath, sheet_name=None)
    except:
        df = pd.read_csv(filepath)
    text = ""
    for sheet_name, sheet in df.items():
        text += f"--- Sheet: {sheet_name} ---\n"
        text += sheet.to_string(index=False) + "\n\n"
    return text

def read_generic(filepath):
    try:
        return textract.process(filepath).decode('utf-8')
    except Exception as e:
        print(f"Cannot extract from {filepath}: {e}")
        return ""

def get_output_filename(base="data", ext=".txt"):
    i = 0
    while True:
        filename = f"{base if i == 0 else base + str(i)}{ext}"
        if not Path(filename).exists():
            return filename
        i += 1

def extract_all_text_from_folder(folder_path):
    all_text = ""
    supported_extensions = {
        ".txt", ".doc", ".docx", ".odt", ".rtf", ".wpd",
        ".ppt", ".pptx", ".odp", ".key",
        ".pdf",
        ".xls", ".xlsx", ".ods",
        ".html", ".htm", ".xml", ".csv", ".log", ".md"
    }

    for root, dirs, files in os.walk(folder_path):
        for file in files:
            ext = Path(file).suffix.lower()
            file_path = os.path.join(root, file)

            print(f"Reading: {file_path}")
            try:
                if ext in {".txt", ".log", ".md", ".html", ".htm", ".xml", ".csv"}:
                    all_text += read_text_file(file_path)
                elif ext == ".docx":
                    all_text += read_docx_file(file_path)
                elif ext == ".pdf":
                    all_text += read_pdf_file(file_path)
                elif ext == ".pptx":
                    all_text += read_pptx_file(file_path)
                elif ext in {".xlsx", ".xls", ".ods"}:
                    all_text += read_spreadsheet(file_path)
                elif ext in supported_extensions:
                    all_text += read_generic(file_path)
                else:
                    print(f"Unsupported format: {file}")
            except Exception as e:
                print(f"Failed to read {file_path}: {e}")
            all_text += "\n\n" + "="*50 + "\n\n"
    
    return all_text

if __name__ == "__main__":
    # 🔧 SET YOUR FOLDER PATH HERE
    folder_path = r"C:\Users\heman\OneDrive\Desktop\adv-rag\Rag-files"

    extracted_text = extract_all_text_from_folder(folder_path)
    output_file = get_output_filename()

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(extracted_text)

    print(f"\n✅ Text data saved to: {output_file}")
    print("Extraction complete! All text has been saved.")